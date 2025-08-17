import os
import re
from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import json

def analyze_content_for_optimal_layout(content: str, slide_title: str = "") -> dict:
    """
    Analyze content to determine the optimal layout and visual elements
    
    Args:
        content: Slide content text
        slide_title: Slide title
    
    Returns:
        Dictionary with layout recommendations and reasoning
    """
    analysis = {
        "layout_type": "content",
        "slide_layout_index": 1,  # Default: Title and Content
        "visual_elements": [],
        "reasoning": "",
        "enhancements": [],
        "content_structure": {}
    }
    
    # Convert to lowercase for analysis
    content_lower = content.lower()
    title_lower = slide_title.lower()
    
    # 1. Check for comparison content
    comparison_keywords = ["vs", "versus", "compared to", "compared with", "difference between", "versus"]
    if any(keyword in content_lower or keyword in title_lower for keyword in comparison_keywords):
        analysis["layout_type"] = "comparison"
        analysis["slide_layout_index"] = 4  # Comparison layout
        analysis["reasoning"] = "Contains comparison keywords - using side-by-side layout"
        analysis["enhancements"].append("side_by_side_layout")
        analysis["content_structure"] = extract_comparison_structure(content)
    
    # 2. Check for numerical data and charts
    number_patterns = [
        r'\d+%',  # Percentages
        r'\d+\s*(percent|%)',  # Percentages with text
        r'increased by \d+',  # Increases
        r'decreased by \d+',  # Decreases
        r'growth of \d+',  # Growth
        r'\d+\s*(million|billion|thousand)',  # Large numbers
    ]
    
    numbers_found = []
    for pattern in number_patterns:
        matches = re.findall(pattern, content_lower)
        numbers_found.extend(matches)
    
    if len(numbers_found) >= 2:
        analysis["layout_type"] = "chart"
        analysis["slide_layout_index"] = 1  # Title and Content with chart
        analysis["reasoning"] = f"Found {len(numbers_found)} numerical data points - adding chart visualization"
        analysis["enhancements"].append("data_visualization")
        analysis["visual_elements"].append("chart")
        analysis["content_structure"] = extract_chart_data(content)
    
    # 3. Check for process/step content
    process_keywords = ["step", "process", "flow", "first", "then", "next", "finally", "stage", "phase"]
    step_patterns = [r'\d+\.', r'step \d+', r'phase \d+']
    
    if any(keyword in content_lower for keyword in process_keywords) or \
       any(re.search(pattern, content_lower) for pattern in step_patterns):
        analysis["layout_type"] = "process"
        analysis["slide_layout_index"] = 1  # Title and Content with SmartArt
        analysis["reasoning"] = "Contains process or step-by-step content - using SmartArt process flow"
        analysis["enhancements"].append("smartart_flow")
        analysis["visual_elements"].append("smartart")
        analysis["content_structure"] = extract_process_structure(content)
    
    # 4. Check for section headers
    section_keywords = ["overview", "introduction", "summary", "conclusion", "key points", "main points"]
    if any(keyword in title_lower for keyword in section_keywords):
        analysis["layout_type"] = "section"
        analysis["slide_layout_index"] = 2  # Section Header layout
        analysis["reasoning"] = "Section header detected - using section layout"
        analysis["enhancements"].append("section_header")
    
    # 5. Check for two-column content
    if content.count('\n') >= 4 and len(content.split('\n')) >= 6:
        # Check if content can be naturally split into two columns
        lines = content.split('\n')
        mid_point = len(lines) // 2
        left_content = '\n'.join(lines[:mid_point])
        right_content = '\n'.join(lines[mid_point:])
        
        if len(left_content) > 20 and len(right_content) > 20:
            analysis["layout_type"] = "two_column"
            analysis["slide_layout_index"] = 3  # Two Content layout
            analysis["reasoning"] = "Content can be split into two balanced columns"
            analysis["enhancements"].append("two_column_layout")
            analysis["content_structure"] = {
                "left_column": left_content,
                "right_column": right_content
            }
    
    # 6. Check for title-only content (impact slides)
    if len(content.strip()) < 50 and any(word in content_lower for word in ["key", "main", "important", "critical"]):
        analysis["layout_type"] = "title_only"
        analysis["slide_layout_index"] = 5  # Title Only layout
        analysis["reasoning"] = "Short, impactful content - using title-only layout for emphasis"
        analysis["enhancements"].append("title_only_emphasis")
    
    # 7. Check for content with caption
    if any(word in content_lower for word in ["diagram", "graph", "chart", "image", "picture", "visual", "illustration"]):
        analysis["layout_type"] = "content_with_caption"
        analysis["slide_layout_index"] = 7  # Content with Caption layout
        analysis["reasoning"] = "Content mentions visual elements - using caption layout"
        analysis["enhancements"].append("caption_layout")
        analysis["visual_elements"].append("image_placeholder")
    
    # 8. Check for blank layout (for custom graphics)
    if len(content.strip()) < 20 and any(word in title_lower for word in ["diagram", "flowchart", "timeline", "mind map"]):
        analysis["layout_type"] = "blank"
        analysis["slide_layout_index"] = 6  # Blank layout
        analysis["reasoning"] = "Minimal content with diagram title - using blank layout for custom graphics"
        analysis["enhancements"].append("blank_custom_graphics")
    
    return analysis

def extract_comparison_structure(content: str) -> dict:
    """Extract comparison data for side-by-side layouts"""
    comparison_data = {
        "left_side": [],
        "right_side": [],
        "comparison_type": "general"
    }
    
    # Split by comparison keywords
    comparison_keywords = ["vs", "versus", "compared to", "compared with"]
    
    for keyword in comparison_keywords:
        if keyword in content.lower():
            parts = content.lower().split(keyword)
            if len(parts) >= 2:
                # Extract bullet points from each side
                left_content = parts[0]
                right_content = parts[1]
                
                # Extract bullet points
                left_points = extract_bullet_points(left_content)
                right_points = extract_bullet_points(right_content)
                
                comparison_data["left_side"] = left_points
                comparison_data["right_side"] = right_points
                break
    
    return comparison_data

def extract_chart_data(content: str) -> dict:
    """Extract numerical data for charts"""
    data = {
        "chart_type": "column",
        "data_points": [],
        "categories": [],
        "values": []
    }
    
    # Extract percentages
    percent_pattern = r'(\d+)%'
    percentages = re.findall(percent_pattern, content)
    if percentages:
        data["chart_type"] = "pie"
        data["values"] = [int(p) for p in percentages]
        data["categories"] = [f"Category {i+1}" for i in range(len(percentages))]
    
    # Extract growth/increase/decrease patterns
    growth_pattern = r'(?:increased|grew|growth)\s+(?:by\s+)?(\d+)'
    growth_matches = re.findall(growth_pattern, content.lower())
    if growth_matches:
        data["chart_type"] = "column"
        data["values"] = [int(g) for g in growth_matches]
        data["categories"] = [f"Growth {i+1}" for i in range(len(growth_matches))]
    
    return data

def extract_process_structure(content: str) -> dict:
    """Extract process steps for SmartArt"""
    process_data = {
        "process_type": "linear",
        "steps": [],
        "total_steps": 0
    }
    
    # Extract numbered steps
    step_pattern = r'(\d+)\.\s*(.+)'
    steps = re.findall(step_pattern, content)
    
    if steps:
        process_data["steps"] = [step[1].strip() for step in steps]
        process_data["total_steps"] = len(steps)
    else:
        # Extract bullet points as steps
        points = extract_bullet_points(content)
        process_data["steps"] = points
        process_data["total_steps"] = len(points)
    
    return process_data

def extract_bullet_points(text: str) -> List[str]:
    """Extract bullet points from text"""
    points = []
    lines = text.split('\n')
    
    for line in lines:
        line = line.strip()
        if line and (line.startswith(('•', '-', '*', '1.', '2.', '3.')) or
                    any(word in line.lower() for word in ['important', 'key', 'main', 'primary'])):
            # Clean up the bullet point
            clean_point = re.sub(r'^[•\-*\d\.\s]+', '', line)
            if clean_point:
                points.append(clean_point)
    
    return points

def create_enhanced_powerpoint(
    title: str,
    slides_data: List[Dict],
    filename: str = "enhanced_presentation.pptx"
) -> str:
    """
    Create an enhanced PowerPoint presentation with smart layout selection
    
    Args:
        title: Main presentation title
        slides_data: List of enhanced slide dictionaries
        filename: Output filename
    
    Returns:
        Path to the created presentation file
    """
    # Create a new presentation
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Set title and subtitle
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = "Smart Layout Enhanced Presentation"
    
    # Process each slide with smart layout selection
    for slide_info in slides_data:
        slide_title = slide_info.get("title", "")
        content = slide_info.get("content", "")
        
        # Analyze content for optimal layout
        layout_analysis = analyze_content_for_optimal_layout(content, slide_title)
        
        # Create slide based on analysis
        slide = create_slide_with_smart_layout(prs, slide_title, content, layout_analysis)
    
    # Save the presentation
    filepath = os.path.join(os.getcwd(), filename)
    prs.save(filepath)
    
    return filepath

def create_slide_with_smart_layout(prs, title: str, content: str, layout_analysis: dict):
    """Create slide using the optimal layout determined by analysis"""
    
    layout_type = layout_analysis["layout_type"]
    slide_layout_index = layout_analysis["slide_layout_index"]
    content_structure = layout_analysis.get("content_structure", {})
    
    try:
        # Get the appropriate layout
        layout = prs.slide_layouts[slide_layout_index]
        slide = prs.slides.add_slide(layout)
        
        # Set title
        slide.shapes.title.text = title
        
        # Create content based on layout type
        if layout_type == "comparison":
            create_comparison_content(slide, content_structure)
        elif layout_type == "chart":
            create_chart_content(slide, content, content_structure)
        elif layout_type == "process":
            create_process_content(slide, content, content_structure)
        elif layout_type == "two_column":
            create_two_column_content(slide, content_structure)
        elif layout_type == "title_only":
            # Title only - no additional content needed
            pass
        elif layout_type == "content_with_caption":
            create_caption_content(slide, content)
        elif layout_type == "blank":
            # Blank layout - ready for custom graphics
            pass
        else:
            # Default content layout
            create_standard_content(slide, content)
        
        return slide
        
    except IndexError:
        # Fallback to standard layout if layout index doesn't exist
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        create_standard_content(slide, content)
        return slide

def create_comparison_content(slide, content_structure: dict):
    """Create side-by-side comparison content"""
    if len(slide.placeholders) >= 3:
        left_placeholder = slide.placeholders[1]
        right_placeholder = slide.placeholders[2]
        
        left_points = content_structure.get("left_side", [])
        right_points = content_structure.get("right_side", [])
        
        add_bullet_points_to_placeholder(left_placeholder, left_points)
        add_bullet_points_to_placeholder(right_placeholder, right_points)

def create_chart_content(slide, content: str, content_structure: dict):
    """Create content with chart visualization"""
    # Add chart if data is available
    chart_data = content_structure.get("data_points", [])
    if chart_data:
        try:
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left=Inches(1),
                top=Inches(2),
                width=Inches(8),
                height=Inches(4)
            ).chart
            
            # Add data to chart
            chart_data_obj = chart.chart_data
            categories = content_structure.get("categories", [])
            values = content_structure.get("values", [])
            
            if categories and values:
                chart_data_obj.categories = categories
                chart_data_obj.add_series("Values", values)
        except:
            pass  # Fallback if chart creation fails
    
    # Add bullet points below chart
    if len(slide.placeholders) >= 2:
        content_placeholder = slide.placeholders[1]
        points = extract_bullet_points(content)
        add_bullet_points_to_placeholder(content_placeholder, points)

def create_process_content(slide, content: str, content_structure: dict):
    """Create process flow content with SmartArt"""
    steps = content_structure.get("steps", [])
    
    if steps:
        try:
            # Create SmartArt process flow
            smartart = slide.shapes.add_smartart(
                smartart_type="process",
                left=Inches(1),
                top=Inches(2),
                width=Inches(8),
                height=Inches(5)
            )
            
            # Add steps to SmartArt
            for i, step in enumerate(steps[:5]):  # Limit to 5 steps
                if i < len(smartart.shapes):
                    smartart.shapes[i].text = step
        except:
            # Fallback to bullet points
            if len(slide.placeholders) >= 2:
                content_placeholder = slide.placeholders[1]
                add_bullet_points_to_placeholder(content_placeholder, steps)

def create_two_column_content(slide, content_structure: dict):
    """Create two-column content layout"""
    if len(slide.placeholders) >= 3:
        left_placeholder = slide.placeholders[1]
        right_placeholder = slide.placeholders[2]
        
        left_content = content_structure.get("left_column", "")
        right_content = content_structure.get("right_column", "")
        
        # Add content to each column
        add_text_to_placeholder(left_placeholder, left_content)
        add_text_to_placeholder(right_placeholder, right_content)

def create_caption_content(slide, content: str):
    """Create content with caption layout"""
    if len(slide.placeholders) >= 2:
        caption_placeholder = slide.placeholders[1]
        text_frame = caption_placeholder.text_frame
        text_frame.text = content[:200] + "..." if len(content) > 200 else content

def create_standard_content(slide, content: str):
    """Create standard bullet point content"""
    if len(slide.placeholders) >= 2:
        content_placeholder = slide.placeholders[1]
        points = extract_bullet_points(content)
        
        if points:
            add_bullet_points_to_placeholder(content_placeholder, points)
        else:
            # Fallback: split content by newlines
            add_text_to_placeholder(content_placeholder, content)

def add_bullet_points_to_placeholder(placeholder, points: List[str]):
    """Add bullet points to a placeholder"""
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0

def add_text_to_placeholder(placeholder, text: str):
    """Add plain text to a placeholder"""
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line.strip()
        p.level = 0

def save_enhanced_plan(slides_data: List[Dict], filename: str = "enhanced_slide_plan.json") -> str:
    """Save enhanced slide plan to file"""
    filepath = os.path.join(os.getcwd(), filename)
    
    with open(filepath, 'w') as f:
        json.dump(slides_data, f, indent=2)
    
    return filepath
