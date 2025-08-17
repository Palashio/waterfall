import os
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

def create_powerpoint(
    title: str,
    slides_data: List[Dict],
    filename: str = "presentation.pptx",
    theme: str = "default"
) -> str:
    """
    Create a PowerPoint presentation from structured data
    
    Args:
        title: Main presentation title
        slides_data: List of slide dictionaries with content
        filename: Output filename
        theme: Presentation theme/style
    
    Returns:
        Path to the created presentation file
    """
    # Create a new presentation
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Set title and subtitle
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"Generated with AI Assistant"
    
    # Process each slide
    for slide_info in slides_data:
        slide_type = slide_info.get("type", "content")
        slide_title = slide_info.get("title", "")
        content = slide_info.get("content", "")
        
        if slide_type == "title":
            # Another title slide
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = slide_title
            slide.placeholders[1].text = content
        elif slide_type == "content":
            # Content slide with bullet points
            content_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_layout)
            
            slide.shapes.title.text = slide_title
            
            # Add content as bullet points
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            # Split content into bullet points
            bullet_points = content.split('\n') if isinstance(content, str) else content
            
            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = point.strip()
                p.level = 0  # Main bullet level
                
        elif slide_type == "section":
            # Section divider slide
            section_layout = prs.slide_layouts[2]  # Section header layout
            slide = prs.slides.add_slide(section_layout)
            slide.shapes.title.text = slide_title
    
    # Save the presentation
    filepath = os.path.join(os.getcwd(), filename)
    prs.save(filepath)
    
    return filepath

def analyze_presentation_structure(
    prompt: str,
    target_slides: int = 5
) -> Dict:
    """
    Analyze a prompt to determine optimal presentation structure
    
    Args:
        prompt: User's presentation request
        target_slides: Target number of slides
    
    Returns:
        Dictionary with presentation structure
    """
    # This is a basic structure - in practice, this would use an LLM
    # to analyze the prompt and determine optimal structure
    
    structure = {
        "title": "Presentation Title",
        "estimated_slides": target_slides,
        "slide_types": ["title", "content", "content", "content", "summary"],
        "key_topics": [],
        "target_audience": "general"
    }
    
    return structure

def save_presentation_plan(
    plan: Dict,
    filename: str = "presentation_plan.json"
) -> str:
    """
    Save the presentation plan to a file
    
    Args:
        plan: Presentation plan dictionary
        filename: Output filename
    
    Returns:
        Path to saved plan file
    """
    filepath = os.path.join(os.getcwd(), filename)
    with open(filepath, 'w') as f:
        json.dump(plan, f, indent=2)
    
    return filepath

def load_presentation_plan(
    filename: str = "presentation_plan.json"
) -> Dict:
    """
    Load a presentation plan from file
    
    Args:
        filename: Plan filename
    
    Returns:
        Loaded plan dictionary
    """
    filepath = os.path.join(os.getcwd(), filename)
    with open(filepath, 'r') as f:
        return json.load(f)
